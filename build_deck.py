# -*- coding: utf-8 -*-
"""Generate SDB AI Insight presentation (.pptx) — clean peach-branded deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette (from product brand) ──────────────────────────────────────────────
PEACH      = RGBColor(0xF4, 0x79, 0x48)
PEACH_DARK = RGBColor(0xC9, 0x4F, 0x28)
PEACH_LT   = RGBColor(0xFF, 0xD9, 0xC2)
PEACH_BG   = RGBColor(0xFD, 0xEC, 0xE0)
INK        = RGBColor(0x2A, 0x24, 0x20)
INK2       = RGBColor(0x6B, 0x63, 0x5E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
CARD       = RGBColor(0xF8, 0xF4, 0xF0)
LINE       = RGBColor(0xEC, 0xE3, 0xDB)

FONT = "Leelawadee UI"   # renders Thai + Latin well on Windows
MONO = "Consolas"

LOGO       = "worker/src/image/softdebutlogo.png"        # colored wordmark (for light bg)
LOGO_WHITE = "worker/src/image/softdebutlogo_white.png"  # white wordmark (for peach bg)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_run_font(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
         sa=6, sb=0, ls=None, font=FONT):
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if sa is not None:
        p.space_after = Pt(sa)
    if sb is not None:
        p.space_before = Pt(sb)
    if ls:
        p.line_spacing = ls
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_run_font(run, font)
    return p


def rect(slide, l, t, w, h, fill, rounded=False, radius=0.12, line=None, lw=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    rect(s, Inches(0.6), Inches(0.58), Inches(0.13), Inches(0.62), PEACH)
    tf = box(s, Inches(0.88), Inches(0.46), Inches(11.8), Inches(1.0))
    if kicker:
        para(tf, kicker, 12, PEACH_DARK, bold=True, sa=2)
        para(tf, title, 27, INK, bold=True)
    else:
        para(tf, title, 28, INK, bold=True, ls=1.0)
    # baseline rule
    rect(s, Inches(0.62), Inches(1.62), Inches(12.1), Pt(1.4), LINE)
    return s


def card(slide, l, t, w, h, fill=CARD, line=LINE):
    return rect(slide, l, t, w, h, fill, rounded=True, radius=0.06, line=line, lw=1.0)


def table(slide, l, t, w, h, data, widths, fs=12.5, hfs=12.5, header=True):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cw in enumerate(widths):
        tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_head = header and r == 0
            if is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = PEACH
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if (r % 2 == 1) else PEACH_BG
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run(); run.text = data[r][c]
            run.font.size = Pt(hfs if is_head else fs)
            run.font.bold = is_head
            run.font.color.rgb = WHITE if is_head else INK
            set_run_font(run, FONT)
    return tbl


def page_no(slide, n):
    tf = box(slide, Inches(12.3), Inches(7.02), Inches(0.8), Inches(0.35))
    para(tf, str(n), 10, INK2, align=PP_ALIGN.RIGHT, sa=0)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
rect(s, 0, 0, Inches(0.35), SH, PEACH)                       # left spine
s.shapes.add_picture(LOGO, Inches(1.0), Inches(2.05), width=Inches(3.0))   # softdebut logo
tf = box(s, Inches(1.0), Inches(3.0), Inches(11.2), Inches(2.6))
para(tf, "SDB AI Insight", 46, INK, bold=True, sa=4)
para(tf, "ระบบติดตามการใช้งาน Claude AI ขององค์กร", 20, INK2, sa=14)
para(tf, "Flow การทำงาน  ·  สถาปัตยกรรม  ·  การทำงานแต่ละหน้า", 14, PEACH_DARK, bold=True)
rect(s, Inches(1.02), Inches(5.4), Inches(4.2), Pt(2.2), PEACH)
tf = box(s, Inches(1.0), Inches(6.7), Inches(9.5), Inches(0.5))
para(tf, "เอกสารนำเสนอระบบ  ·  อัปเดต 2026-05-20", 12, INK2, sa=0)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("สิ่งที่จะเล่าวันนี้", kicker="AGENDA")
cw, ch, top = Inches(5.85), Inches(4.2), Inches(2.1)
gap = Inches(0.45)
left1 = Inches(0.62)
left2 = left1 + cw + gap
for (lx, num, title, desc) in [
    (left1, "01", "Flow การทำงานของระบบ",
     ["• ภาพรวมและปัญหาที่แก้", "• สถาปัตยกรรม 3 ชั้น", "• การเก็บข้อมูล 1 call", "• การตรวจ client + identity", "• การคิดเงิน · ฐานข้อมูล · auth"]),
    (left2, "02", "แต่ละหน้าทำอะไร",
     ["• Dashboard · Accounts", "• Account Detail · Analytics", "• Monitoring · Data Sources", "• Identity · Reports · Settings", "• Clear Data · Insights"]),
]:
    card(s, lx, top, cw, ch)
    rect(s, lx + Inches(0.35), top + Inches(0.35), Inches(1.1), Inches(0.95), PEACH_BG, rounded=True, radius=0.2)
    t = box(s, lx + Inches(0.35), top + Inches(0.35), Inches(1.1), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(t, num, 30, PEACH, bold=True, align=PP_ALIGN.CENTER, sa=0)
    t = box(s, lx + Inches(0.35), top + Inches(1.55), cw - Inches(0.7), Inches(2.4))
    para(t, title, 22, INK, bold=True, sa=10)
    for d in desc:
        para(t, d, 14, INK2, sa=6, ls=1.05)
page_no(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# Section divider helper
# ════════════════════════════════════════════════════════════════════════════
def section(num, kicker, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PEACH)
    t = box(s, Inches(7.0), Inches(0.4), Inches(6.0), Inches(3.0))
    para(t, num, 140, PEACH_DARK, bold=True, align=PP_ALIGN.RIGHT, sa=0)
    rect(s, Inches(0.9), Inches(2.95), Inches(0.9), Pt(3), WHITE)
    t = box(s, Inches(0.9), Inches(3.2), Inches(10.5), Inches(2.6))
    para(t, kicker, 15, PEACH_LT, bold=True, sa=8)
    para(t, title, 40, WHITE, bold=True, sa=12)
    para(t, subtitle, 17, PEACH_LT, ls=1.1)
    return s


# SLIDE 3 — PART 1 divider
section("1", "ส่วนที่ 1", "Flow การทำงานของระบบ",
        "ภาพรวม · สถาปัตยกรรม 3 ชั้น · การเก็บข้อมูล · การตรวจตัวตน · การคิดเงิน")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is it
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("SDB AI Insight คืออะไร", kicker="ภาพรวม")
tf = box(s, Inches(0.62), Inches(2.0), Inches(12.0), Inches(1.5))
para(tf, "ระบบ monitoring การใช้งาน Claude AI ขององค์กร — ดักจับทุก prompt/response "
         "จากทุก client แล้วสรุปเป็นค่าใช้จ่าย (USD), token usage และพฤติกรรมการใช้งาน บน dashboard",
     17, INK, ls=1.2)
card(s, Inches(0.62), Inches(3.55), Inches(12.1), Inches(3.1), fill=PEACH_BG, line=PEACH_LT)
tf = box(s, Inches(1.0), Inches(3.85), Inches(11.4), Inches(2.6))
para(tf, "ปัญหาที่แก้  —  องค์กรจ่ายค่า Claude แต่มองไม่เห็นว่า:", 16, PEACH_DARK, bold=True, sa=12)
for d in [
    "ใครใช้บ้าง · ใช้ไปเท่าไหร่ · model ไหนกินเงินมากที่สุด",
    "ใช้ผ่านช่องทางไหน (Desktop, Claude.ai, Cowork, Code, CLI, VSCode, API)",
    "ค่าใช้จ่ายพุ่งสูงผิดปกติในช่วงเวลาใด",
]:
    para(tf, "●  " + d, 15, INK, sa=9, ls=1.1)
page_no(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Tech stack
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("Tech Stack", kicker="เทคโนโลยีหลัก 4 ส่วน")
items = [
    ("mitmproxy", "Python proxy", "ดักจับ traffic + วิเคราะห์"),
    ("Cloudflare Workers", "Edge runtime", "รับข้อมูล + render dashboard"),
    ("Cloudflare D1", "SQLite database", "เก็บ log ทั้งหมด"),
    ("Logto", "OIDC + PKCE", "ยืนยันตัวตนเข้า dashboard"),
]
cw = Inches(2.86); gap = Inches(0.27); top = Inches(2.35); ch = Inches(3.6)
lx = Inches(0.62)
for (name, tag, desc) in items:
    card(s, lx, top, cw, ch)
    rect(s, lx, top, cw, Inches(0.16), PEACH, rounded=False)
    t = box(s, lx + Inches(0.28), top + Inches(0.55), cw - Inches(0.56), ch - Inches(0.8))
    para(t, name, 17, INK, bold=True, sa=6, ls=1.0)
    para(t, tag, 12, PEACH_DARK, bold=True, sa=14)
    para(t, desc, 14, INK2, ls=1.15)
    lx = lx + cw + gap
page_no(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture 3 layers
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("สถาปัตยกรรม 3 ชั้น", kicker="ARCHITECTURE")
cols = [
    ("1", "CLIENTS", ["Claude Desktop", "Claude.ai web", "Cowork · Code tab", "Code CLI · VSCode", "API SDK"]),
    ("2", "PROXY  (mitmproxy + addon.py)", ["ดักจับ traffic (TLS MITM)", "ตรวจ client · ดึง prompt", "นับ token · คิดเงิน", "ตรวจ email เจ้าของ", "→ JSONL + POST /log"]),
    ("3", "WORKER + D1", ["รับ log ที่ /log", "เติม identity ที่ขาด", "เก็บลง D1", "Logto login", "render dashboard"]),
]
cw = Inches(3.86); gap = Inches(0.26); top = Inches(2.15); ch = Inches(3.85)
lx = Inches(0.62)
for (num, title, rows) in cols:
    card(s, lx, top, cw, ch)
    rect(s, lx + Inches(0.3), top + Inches(0.3), Inches(0.62), Inches(0.62), PEACH, rounded=True, radius=0.25)
    t = box(s, lx + Inches(0.3), top + Inches(0.3), Inches(0.62), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    para(t, num, 20, WHITE, bold=True, align=PP_ALIGN.CENTER, sa=0)
    t = box(s, lx + Inches(1.05), top + Inches(0.34), cw - Inches(1.3), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    para(t, title, 13.5, INK, bold=True, ls=0.95, sa=0)
    t = box(s, lx + Inches(0.32), top + Inches(1.2), cw - Inches(0.6), ch - Inches(1.4))
    for r in rows:
        para(t, "•  " + r, 13.5, INK2, sa=8, ls=1.05)
    lx = lx + cw + gap
tf = box(s, Inches(0.62), Inches(6.25), Inches(12), Inches(0.5))
para(tf, "ทุกเครื่อง client ตั้ง system proxy → 127.0.0.1:8080  ·  proxy เก็บข้อมูล / worker เก็บถาวร+แสดงผล แยกหน้าที่กัน",
     12.5, PEACH_DARK, bold=True, align=PP_ALIGN.CENTER)
page_no(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Flow of one call (5 steps)
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("Flow การเก็บ 1 call แบบครบวงจร", kicker="DATA CAPTURE")
steps = [
    ("1", "ผู้ใช้พิมพ์ prompt", "traffic วิ่งผ่าน mitmproxy อัตโนมัติ (เพราะตั้ง system proxy ไว้)"),
    ("2", "addon.py วิเคราะห์", "detect client · ดึง prompt จริง · นับ token · คิดเงิน · หา email เจ้าของ"),
    ("3", "ตรวจ email filter", "ผ่าน → เขียนไฟล์ JSONL ในเครื่อง (backup) + POST /log แบบ fire-and-forget"),
    ("4", "Worker /log", "ตรวจ X-Api-Key · เติม email จาก ip_identity ถ้าว่าง · INSERT api_logs · upsert ตัวตน"),
    ("5", "ขึ้น Dashboard", "ข้อมูลปรากฏทันที — dashboard query สดจาก D1 ทุกครั้งที่เปิด"),
]
top = Inches(2.05); rh = Inches(0.92); lx = Inches(0.62)
for (num, title, desc) in steps:
    card(s, lx, top, Inches(12.1), rh - Inches(0.12))
    rect(s, lx + Inches(0.2), top + Inches(0.14), Inches(0.55), Inches(0.55), PEACH, rounded=True, radius=0.25)
    t = box(s, lx + Inches(0.2), top + Inches(0.14), Inches(0.55), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(t, num, 18, WHITE, bold=True, align=PP_ALIGN.CENTER, sa=0)
    t = box(s, lx + Inches(1.0), top + Inches(0.06), Inches(2.7), rh - Inches(0.2), anchor=MSO_ANCHOR.MIDDLE)
    para(t, title, 15.5, INK, bold=True, sa=0, ls=1.0)
    t = box(s, lx + Inches(3.8), top + Inches(0.06), Inches(8.1), rh - Inches(0.2), anchor=MSO_ANCHOR.MIDDLE)
    para(t, desc, 13.5, INK2, sa=0, ls=1.05)
    top = top + rh
page_no(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Proxy classes
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("ชั้น Proxy: addon.py ทำอะไร", kicker="CAPTURE LAYER")
data = [
    ["Class", "ดักอะไร / หน้าที่"],
    ["ClaudeAPIMonitor", "api.anthropic.com/v1/messages (API key, Code CLI/VSCode, Cowork, Code tab)"],
    ["ClaudeDesktopMonitor", "claude.ai .../completion (Claude Desktop & web chat)"],
    ["ClaudeBridgeMonitor", "bridge.claudeusercontent.com WebSocket (Claude Code OAuth login)"],
    ["ClaudeAccountSniffer", "อ่าน email จาก claude.ai — ทำงานก่อนเสมอ"],
    ["ToolSchemaFixer", "แก้ tool schema ที่ Anthropic API ปฏิเสธ ให้คำขอผ่าน"],
]
table(s, Inches(0.62), Inches(2.05), Inches(12.1), Inches(3.5), data,
      widths=[Inches(3.4), Inches(8.7)], fs=13, hfs=13)
tf = box(s, Inches(0.62), Inches(5.95), Inches(12), Inches(0.6))
para(tf, "+ Discovery classes (debug): ConnectionLogger, DesktopDiscovery, BridgeDiscovery — บันทึก endpoint/frame ที่ยังไม่รู้จัก",
     12.5, INK2, align=PP_ALIGN.CENTER)
page_no(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Client detection
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("การตรวจจับว่ามาจาก Client ใด", kicker="CLIENT DETECTION")
tf = box(s, Inches(0.62), Inches(1.95), Inches(12), Inches(0.5))
para(tf, "ชั้น 1 — Headers  (User-Agent / anthropic-client-name / x-app / x-client-context)", 15, PEACH_DARK, bold=True)
data = [
    ["เงื่อนไข", "ผลลัพธ์ (client tag)"],
    ["claude-code + electron", "claude-desktop-code"],
    ["claude-code + vscode", "claude-code-vscode"],
    ["claude-code เพียวๆ", "claude-code-cli"],
    ["electron / anthropic ใน UA", "claude-desktop"],
    ["ไม่เข้าเงื่อนไขใด", "api"],
]
table(s, Inches(0.62), Inches(2.45), Inches(7.4), Inches(3.0), data,
      widths=[Inches(4.3), Inches(3.1)], fs=13, hfs=12.5)
card(s, Inches(8.35), Inches(2.45), Inches(4.37), Inches(3.0), fill=PEACH_BG, line=PEACH_LT)
t = box(s, Inches(8.65), Inches(2.7), Inches(3.8), Inches(2.6))
para(t, "ชั้น 2 — Body override", 14.5, PEACH_DARK, bold=True, sa=10)
para(t, "เมื่อ header กำกวม:", 13, INK, sa=10)
para(t, "● body มี mcp__cowork__*", 13, INK, sa=2, ls=1.05)
para(t, "   → claude-desktop-cowork (ชนะเสมอ)", 12.5, INK2, sa=10, ls=1.05)
para(t, "● body มี Code tools (Bash/Read/...)", 13, INK, sa=2, ls=1.05)
para(t, "   แต่ header = api → claude-code-cli", 12.5, INK2, ls=1.05)
page_no(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Identity 4 layers
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("การระบุตัวตนแบบ 4 ชั้น", kicker="IDENTITY · 4 LAYERS")
tf = box(s, Inches(0.62), Inches(1.92), Inches(12.1), Inches(0.7))
para(tf, "ปัญหา: บน proxy กลาง บาง log มี email ว่าง (25–64%) → เก็บตัวตนกระจาย 4 ชั้น แต่ละชั้นรอดแยกกันได้",
     14, INK, ls=1.15)
data = [
    ["Layer", "เก็บที่ไหน", "บทบาท"],
    ["L1  _ACCOUNT_BY_IP", "Proxy (in-memory)", "hot read path — map IP → email ตอน log"],
    ["L2  account_slots.json", "Proxy (ไฟล์)", "สำเนาของ L1 — รอด proxy restart"],
    ["L3  ip_identity", "Worker D1", "source of truth ระยะยาว — รอดทุกอย่าง"],
    ["L4  api_logs.client_ip", "Worker D1", "audit ทุก row + fallback แสดงผล"],
]
table(s, Inches(0.62), Inches(2.7), Inches(12.1), Inches(2.7), data,
      widths=[Inches(3.5), Inches(3.0), Inches(5.6)], fs=12.5, hfs=12.5)
tf = box(s, Inches(0.62), Inches(5.6), Inches(12.1), Inches(1.0))
para(tf, "ได้ email จาก 3 แหล่ง: HTTP sniffer (claude.ai) · JWT decode (Bearer token) · Bridge connect (WebSocket)", 13, INK2, sa=6, ls=1.1)
para(tf, "หา email ไม่ได้เลย → แสดงเป็น ip:10.10.84.42 แทน (ยัง audit ได้ว่ามาจากเครื่องไหน)", 13, PEACH_DARK, bold=True)
page_no(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Pricing
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("การคำนวณค่าใช้จ่าย", kicker="PRICING (USD / 1M tokens)")
data = [
    ["Model tier", "Input", "Output", "Cache Read", "Cache Write"],
    ["Opus", "$15", "$75", "$1.50", "$18.75"],
    ["Sonnet", "$3", "$15", "$0.30", "$3.75"],
    ["Haiku", "$0.80", "$4", "$0.08", "$1.00"],
]
table(s, Inches(0.62), Inches(2.15), Inches(12.1), Inches(2.3), data,
      widths=[Inches(3.3), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)], fs=14, hfs=13.5)
card(s, Inches(0.62), Inches(4.85), Inches(12.1), Inches(1.7), fill=PEACH_BG, line=PEACH_LT)
t = box(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(1.3), anchor=MSO_ANCHOR.MIDDLE)
para(t, "cost = (input×Pin + output×Pout + cacheRead×Pcr + cacheWrite×Pcw) ÷ 1,000,000",
     16, INK, bold=True, font=MONO, sa=8, align=PP_ALIGN.CENTER)
para(t, "แยกนับ cache เพราะราคาต่างกันมาก — cache read ถูกกว่า input ปกติ ~10 เท่า",
     13, INK2, align=PP_ALIGN.CENTER)
page_no(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Worker endpoints + auth
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("Worker Endpoints & Auth", kicker="BACKEND")
data = [
    ["Path", "Method", "หน้าที่"],
    ["/log", "POST", "รับ log จาก proxy (ตรวจ X-Api-Key)"],
    ["/health", "GET", "health check"],
    ["/login", "GET", "redirect ไป Logto"],
    ["/  (?code=)", "GET", "OAuth callback"],
    ["/logout", "GET", "จบ session"],
]
table(s, Inches(0.62), Inches(2.05), Inches(7.4), Inches(3.1), data,
      widths=[Inches(2.5), Inches(1.5), Inches(3.4)], fs=12.5, hfs=12.5)
card(s, Inches(8.35), Inches(2.05), Inches(4.37), Inches(3.1), fill=PEACH_BG, line=PEACH_LT)
t = box(s, Inches(8.65), Inches(2.3), Inches(3.8), Inches(2.7))
para(t, "Auth = Logto OIDC + PKCE", 14.5, PEACH_DARK, bold=True, sa=10)
for d in ["/login → Logto authorize", "callback (?code=) → แลก token",
          "verify id_token (JWKS)", "สร้าง session (cookie sid, 7 วัน)", "→ เข้า dashboard"]:
    para(t, "→  " + d, 13, INK, sa=7, ls=1.05)
tf = box(s, Inches(0.62), Inches(5.5), Inches(12), Inches(0.6))
para(tf, "หน้า dashboard ทั้งหมดต้องผ่าน requireUser gate ก่อนเสมอ", 13, INK2, align=PP_ALIGN.CENTER)
page_no(s, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Database
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("ฐานข้อมูล Cloudflare D1", kicker="DATABASE")
data = [
    ["Table", "เก็บอะไร"],
    ["api_logs", "log ทุก call — ts, client, account_email, client_ip, model, prompt, tokens, cost_usd"],
    ["ip_identity", "map IP → email (Layer 3) — เติม email ให้ log ที่ระบุตัวไม่ได้"],
    ["sessions", "session ผู้ใช้ที่ login (sub, email, expires_at, id_token)"],
    ["oauth_state", "state + PKCE verifier ระหว่าง OAuth flow"],
    ["app_settings", "config แบบ key/value (ingest_key, notify flags)"],
]
table(s, Inches(0.62), Inches(2.05), Inches(12.1), Inches(3.5), data,
      widths=[Inches(2.7), Inches(9.4)], fs=12.5, hfs=12.5)
tf = box(s, Inches(0.62), Inches(5.95), Inches(12), Inches(0.5))
para(tf, "ทุกเวลาคำนวณตามเขต Asia/Bangkok (UTC+7) · เก็บค่าเวลาเป็น ms epoch", 12.5, INK2, align=PP_ALIGN.CENTER)
page_no(s, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Coverage
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("ช่องทาง Claude ที่ระบบครอบคลุม", kicker="COVERAGE")
data = [
    ["แหล่งที่มา", "ช่องทาง", "client tag"],
    ["Claude Code CLI (API key)", "api.anthropic.com/v1/messages", "claude-code-cli"],
    ["Claude Code CLI (OAuth)", "bridge.claudeusercontent.com WS", "claude-code-cli"],
    ["Claude Code VSCode", "api.anthropic.com/v1/messages", "claude-code-vscode"],
    ["Claude Desktop / web chat", "claude.ai .../completion", "claude-desktop"],
    ["Cowork (Desktop)", "/v1/messages?beta=true", "claude-desktop-cowork"],
    ["Code tab (Desktop)", "/v1/messages?beta=true", "claude-desktop-code"],
    ["Claude API SDK", "api.anthropic.com/v1/messages", "api"],
]
table(s, Inches(0.62), Inches(2.0), Inches(12.1), Inches(4.4), data,
      widths=[Inches(4.0), Inches(4.7), Inches(3.4)], fs=12, hfs=12)
page_no(s, 14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Limitations
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("ข้อจำกัดที่ทราบ", kicker="LIMITATIONS")
tf = box(s, Inches(0.7), Inches(2.15), Inches(12), Inches(4.5))
for d in [
    "Mobile apps — อยู่คนละเครื่อง ต้อง MITM ที่ network layer (router)",
    "HTTP/3 (QUIC) — mitmproxy ดักได้แค่ TCP (ปัจจุบันยังไม่กระทบ)",
    "API key users — ไม่มี email ใน traffic → ระบุได้แค่ IP (fundamental)",
    "NAT / shared IP — หลายคนใช้ IP เดียวกัน → แยกตัวตนไม่ออก",
    "First call ของ IP ใหม่ — ต้องรอ sniff/JWT ก่อน 1 call จึงได้ email",
]:
    para(tf, "●   " + d, 16, INK, sa=16, ls=1.15)
page_no(s, 15)


# SLIDE 16 — PART 2 divider
section("2", "ส่วนที่ 2", "แต่ละหน้าทำอะไร ทำงานอย่างไร",
        "เมนูหลัก: Dashboard · Analytics · Accounts  |  หน้าอื่นเข้าผ่าน URL ได้")


# ── page-with-bullets helper ───────────────────────────────────────────────────
def page_bullets(title, kicker, route, intro, bullets, n):
    s = content_slide(title, kicker=kicker)
    tf = box(s, Inches(0.62), Inches(1.78), Inches(12), Inches(0.45))
    para(tf, "route: " + route, 12.5, PEACH_DARK, bold=True, font=MONO)
    tf = box(s, Inches(0.62), Inches(2.25), Inches(12.1), Inches(0.9))
    para(tf, intro, 15.5, INK, ls=1.2)
    tf = box(s, Inches(0.7), Inches(3.25), Inches(12), Inches(3.5))
    for b in bullets:
        para(tf, "●   " + b, 14.5, INK2, sa=11, ls=1.12)
    page_no(s, n)
    return s


# SLIDE 17 — Dashboard
page_bullets("หน้า Dashboard", "หน้าหลัก", "/",
    "สรุปการใช้งานทั้งหมดพร้อมตัวกรองครบถ้วน",
    ["Filter bar — Period (Daily/Monthly/Yearly), Date from–to, Model, Account, Client + Export CSV / Apply",
     "Stat cards 6 ใบ — calls, token in/out, cache read/write, cost รวม",
     "Breakdown 3 การ์ด — By Model · By Account · By Client (เรียงตาม cost)",
     "Recent API Calls — ตารางทุก call: เวลา, client, account, model, prompt, token, cost",
     "เลือก rows/หน้า (10–All) + แบ่งหน้า · คลิก prompt เปิด modal ดูข้อความเต็ม"], 17)


# SLIDE 18 — Accounts + Account Detail
s = content_slide("หน้า Accounts & Account Detail", kicker="บัญชีผู้ใช้")
left, top, cw, ch = Inches(0.62), Inches(2.05), Inches(5.9), Inches(4.5)
card(s, left, top, cw, ch)
rect(s, left, top, cw, Inches(0.16), PEACH)
t = box(s, left + Inches(0.3), top + Inches(0.45), cw - Inches(0.6), ch - Inches(0.7))
para(t, "Accounts  ·  /accounts", 16, INK, bold=True, sa=10)
for b in ["รายชื่อบัญชีทั้งหมดที่เชื่อม Claude", "Period filter: 7d / 30d / 90d / all",
          "KPI: total/active accounts, spend, calls", "ตาราง: avatar, email, calls, tokens, model, cost",
          "สถานะ: live (<1ชม) · idle (<7วัน) · cold", "คลิกแถว → Account Detail"]:
    para(t, "•  " + b, 13.5, INK2, sa=8, ls=1.08)
left2 = left + cw + Inches(0.3)
card(s, left2, top, cw, ch)
rect(s, left2, top, cw, Inches(0.16), PEACH)
t = box(s, left2 + Inches(0.3), top + Inches(0.45), cw - Inches(0.6), ch - Inches(0.7))
para(t, "Account Detail  ·  /account?identity=", 16, INK, bold=True, sa=10)
for b in ["เจาะลึกบัญชีเดียว (รับทั้ง email และ ip:...)", "Stat cards + first/last seen + สถานะ",
          "Cost over time (30 วัน)", "By Model / By Client ของคนนี้",
          "Top 5 prompts ที่ใช้บ่อย · Token usage", "Activity heatmap 7×24 · 50 prompts ล่าสุด"]:
    para(t, "•  " + b, 13.5, INK2, sa=8, ls=1.08)
page_no(s, 18)


# SLIDE 19 — Analytics (the 5 charts — important)
s = content_slide("หน้า Analytics", kicker="แนวโน้มเชิงลึก · /analytics")
tf = box(s, Inches(0.62), Inches(1.8), Inches(12), Inches(0.5))
para(tf, "ดู trend ค่าใช้จ่ายและ token ตามช่วงเวลา (7d / 30d / 90d) — มีกราฟ 5 ตัว", 14.5, INK, ls=1.1)
data = [
    ["กราฟ", "คืออะไร"],
    ["Cost over time", "ราคาที่จ่ายไป แยกตามช่วงเวลา (เช่น 2 วัน: 18/05 และ 19/05) — 1 จุด = 1 วัน"],
    ["Cost by model", "เทียบ trend ค่าใช้จ่ายของแต่ละ model (top 7) ว่าตัวไหนกินเงินมากสุด"],
    ["Token mix over time", "จำนวน token แบ่งตามเวลา — 'mix' = รวม 4 ตัวซ้อนกัน: cache read/write + input/output"],
    ["Calls timeline", "จำนวนครั้งที่เรียกใช้ (requests) ต่อวัน — นับเป็นจำนวน call ไม่ใช่ยอดเงิน"],
    ["Activity heatmap", "ความถี่การ prompt แต่ละช่วงเวลาของวัน — ตาราง 7×24 ยิ่งเข้มยิ่งใช้บ่อย"],
]
table(s, Inches(0.62), Inches(2.35), Inches(12.1), Inches(3.4), data,
      widths=[Inches(2.9), Inches(9.2)], fs=12.5, hfs=12.5)
card(s, Inches(0.62), Inches(5.95), Inches(12.1), Inches(0.85), fill=PEACH_BG, line=PEACH_LT)
t = box(s, Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
para(t, "⚠  Calls timeline = จำนวน call/วัน (count) ไม่ใช่เงิน — ถ้าจะสื่อ 'ค่าใช้จ่ายรวมต่อวัน' ให้ดู Cost over time แทน",
     12.5, PEACH_DARK, bold=True, align=PP_ALIGN.CENTER)
page_no(s, 19)


# SLIDE 20 — Monitoring + Data Sources
s = content_slide("หน้า Monitoring & Data Sources", kicker="ระบบ & แหล่งข้อมูล")
left, top, cw, ch = Inches(0.62), Inches(2.05), Inches(5.9), Inches(4.5)
card(s, left, top, cw, ch)
rect(s, left, top, cw, Inches(0.16), PEACH)
t = box(s, left + Inches(0.3), top + Inches(0.45), cw - Inches(0.6), ch - Inches(0.7))
para(t, "Monitoring  ·  /monitoring", 16, INK, bold=True, sa=10)
for b in ["สถานะ real-time (ย้อนหลัง 24 ชม.)", "KPI: calls, error rate, active sessions, status",
          "API Health (error rate/ชม.)", "Active sessions · Throughput (calls/min)",
          "Top errors by group", "error = output=0 และ cost=0 (ไม่มีคำตอบ)"]:
    para(t, "•  " + b, 13.5, INK2, sa=8, ls=1.08)
left2 = left + cw + Inches(0.3)
card(s, left2, top, cw, ch)
rect(s, left2, top, cw, Inches(0.16), PEACH)
t = box(s, left2 + Inches(0.3), top + Inches(0.45), cw - Inches(0.6), ch - Inches(0.7))
para(t, "Data Sources  ·  /data-sources", 16, INK, bold=True, sa=10)
for b in ["จัดการ proxy ที่ป้อนข้อมูลเข้าระบบ", "Ingest endpoint (/log) + ingest key (mask)",
          "สถิติ D1: rows รวม, rows วันนี้,", "    ช่วงเวลา oldest–newest, ขนาดข้อมูล",
          "คำแนะนำการตั้งค่า proxy ให้ client"]:
    para(t, "•  " + b, 13.5, INK2, sa=8, ls=1.08)
page_no(s, 20)


# SLIDE 21 — Other pages
s = content_slide("หน้าอื่นๆ ในระบบ", kicker="MORE PAGES")
pages = [
    ("Identity  /identity", "ตาราง map IP ↔ email ปัจจุบัน (Layer 3) — เติม email ให้ log ที่ระบุตัวไม่ได้"),
    ("Reports  /reports", "Export CSV: Today / Month / Year / All + custom report (เลือกช่วง/columns)"),
    ("Settings  /settings", "Profile · Rotate ingest key · Notifications (email/anomaly/budget) · About"),
    ("Clear Data  /clear-data", "ลบข้อมูล: All / ตามช่วงวันที่ / ตาม filter / Sessions — ยืนยันจำนวนที่ลบ"),
    ("Insights  /insights", "⏳ Incoming — AI-driven findings (อยู่ระหว่างพัฒนา)"),
]
top = Inches(2.05); rh = Inches(0.95); lx = Inches(0.62)
for (name, desc) in pages:
    card(s, lx, top, Inches(12.1), rh - Inches(0.14))
    t = box(s, lx + Inches(0.32), top + Inches(0.04), Inches(3.5), rh - Inches(0.2), anchor=MSO_ANCHOR.MIDDLE)
    para(t, name, 14.5, INK, bold=True, sa=0, font=MONO, ls=1.0)
    t = box(s, lx + Inches(3.95), top + Inches(0.04), Inches(7.9), rh - Inches(0.2), anchor=MSO_ANCHOR.MIDDLE)
    para(t, desc, 13, INK2, sa=0, ls=1.05)
    top = top + rh
page_no(s, 21)


# SLIDE 22 — Key takeaways
s = content_slide("สรุปประเด็นสำคัญ", kicker="KEY TAKEAWAYS")
tf = box(s, Inches(0.7), Inches(2.05), Inches(12), Inches(4.6))
for d in [
    "ดักทุกช่องทาง — Desktop, web, Cowork, Code tab, CLI, VSCode, API ในที่เดียว",
    "คิดเงินแม่นยำ — แยก input/output/cache ตามราคาจริงของแต่ละ model",
    "ระบุตัวตนทน fail — 4-layer identity รอด proxy/worker restart",
    "ไม่ block ผู้ใช้ — log แบบ fire-and-forget + backup JSONL ในเครื่อง",
    "Dashboard ครบ — ภาพรวม, รายคน, trend, real-time monitoring, export",
]:
    para(tf, "✓   " + d, 16.5, INK, bold=False, sa=15, ls=1.15)
page_no(s, 22)


# SLIDE 23 — Closing
s = prs.slides.add_slide(BLANK)
set_bg(s, PEACH)
_lw = Inches(3.9)
s.shapes.add_picture(LOGO_WHITE, (SW - _lw) // 2, Inches(2.5), width=_lw)   # white softdebut logo
t = box(s, Inches(1.0), Inches(3.95), Inches(11.3), Inches(1.6), anchor=MSO_ANCHOR.TOP)
para(t, "ขอบคุณครับ", 40, WHITE, bold=True, align=PP_ALIGN.CENTER, sa=8)
para(t, "SDB AI Insight — ติดตามการใช้งาน Claude AI ได้ครบในที่เดียว", 17, PEACH_LT, align=PP_ALIGN.CENTER)

prs.save("SDB-AI-Insight.pptx")
print("OK saved SDB-AI-Insight.pptx :", len(prs.slides._sldIdLst), "slides")
